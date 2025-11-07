# controller/main_controller.py
import getpass
import os
import pandas as pd
from PyQt5.QtCore import QObject, pyqtSignal, Qt
from PyQt5.QtWidgets import QMessageBox
from datetime import datetime, timedelta
from controllers.field_controller import FieldController
from data.data_loader import DataLoader
from logic.activity_data import build_activities_dataframe
from logic.operative_capacity_manager import OperativeCapacityManager
from logic.reports.rig_report import RigReport
from logic.reports.mi_swaco_report import MISwacoReport
from logic.reports.completions import CompletionsReport
from logic.reports.bits_drilling_remedial import BitsDrillingTRemedialReport
from logic.reports.surface_systems import SurfaceSystemsReport
from logic.reports.well_services_report import WellServicesReport
from logic.reports.testing_fluid_analysis import TestingFluidAnalysisReport
from logic.reports.wireline_report import WirelineReport
from logic.reports.services import ServicesReport
from logic.reports.artificial_lift_report import ArtificialLiftReport
from logic.reports.integrated_services_report import IntegratedServicesReport
from logic.reports.environment_report import EnvironmentReport    
from logic.reports.tubulars_report import TubularsReport
from logic.reports.tanks_and_trunks_report import TanksAndTrunksReport
from services.capex_config_service import CapexConfigService
from utils.file_manager import (
    get_catalog_dir, get_forecast_services_path_file, get_forecasted_plan_path, get_operative_capacity_path, get_plan_path, get_budget_opex_path, get_planning_cost_path
)
from pptx import Presentation
from pptx.util import Inches
from utils.export_ppt import add_slide_to_presentation
from PyQt5.QtWidgets import QDialog

from views.capex_config_view import CapexConfigDialog
from views.field_views.avg_days_dialog import AvgDaysDialog
from views.office_planning_view import OfficePlanningView
from views.services_forecast_path_view import ServicesForecastPathView

from utils.comments import load_comments, save_comment, load_field_line_comments
from logic.opex_data_manager import OpexDataManager

# ... (importaciones de nuevas vistas)
from views.capex_config_view import CapexConfigDialog
from views.field_views.avg_days_dialog import AvgDaysDialog
# ...
from views.tubulars_config import TubularsConfigDialog
# --- AÑADIDO PARA MI SWACO Y COMPLETIONS ---
from views.mi_swaco_config_view import MISwacoConfigDialog
from views.completions_config_view import CompletionsConfigDialog

from views.opex_editor import OpexEditorWindow
from logic.plan_actividades import PlanAnualActividades

from views.services_resumen_dialog import ServicesResumenDialog
from views.well_selector import WellSelectorDialog
from datetime import datetime, timedelta
from utils.file_manager import get_output_path_for_pptx, get_selected_services_wells_path, get_planned_activities_catalog_path

# --- Clases de Respaldo (Fallback) ---
# Se usan si los archivos de configuración de oficina no se encuentran,
# para evitar que la aplicación falle al iniciar para usuarios de campo.

class _EmptyCapacityManager:
    """Gestor de capacidad operativa vacío de respaldo."""
    def __init__(self, *args, **kwargs):
        self.df = pd.DataFrame()
    def get_total_tentative_opex_wells(self): return 0
    def update_value(self, *args, **kwargs): pass
    def save(self, *args, **kwargs): pass
    def export_to(self, *args, **kwargs): pass

class _EmptyPlanActividades:
    """Plan de actividades vacío de respaldo."""
    def __init__(self, *args, **kwargs):
        self.plan_df = pd.DataFrame()
        self.plan_path = ""
        self.sheet_name = ""

class _EmptyOpexDataManager:
    """Gestor de datos OPEX vacío de respaldo."""
    def __init__(self, *args, **kwargs): pass
    def load_opex_data(self): return pd.DataFrame()
    def set_opex_data(self, df): pass
    def save_opex_to_excel(self): pass

class MainController(QObject):    
    dataUpdated = pyqtSignal()
    def __init__(self):
        super().__init__()
        # -------------------------
        # 🧠 Inicialización general
        # -------------------------
        self.data_loader = DataLoader()
        try:
            self.comments_df = load_comments()
        except (FileNotFoundError, pd.errors.EmptyDataError) as e:
            print(f"⚠️ ADVERTENCIA: No se pudo cargar el archivo de comentarios de oficina: {e}. Se usará un DataFrame vacío.")
            self.comments_df = pd.DataFrame(columns=["Report Title", "Mes", "Fecha", "Usuario", "Comentario"])
        
        try:
            self.field_comments_df = load_field_line_comments()
        except (FileNotFoundError, pd.errors.EmptyDataError) as e:
            print(f"⚠️ ADVERTENCIA: No se pudo cargar el archivo de comentarios de campo: {e}. Se usará un DataFrame vacío.")
            # Las columnas se infieren más adelante si es necesario, un DF vacío es seguro.
            self.field_comments_df = pd.DataFrame()
        self.year_actual = datetime.now().year
        self.capex_config_service = CapexConfigService()
        # -------------------------
        # 📁 Rutas de archivos
        # -------------------------
        self.catalog_dir = get_catalog_dir()
        self.operative_capacity_file = get_operative_capacity_path()
        self.plan_path = get_plan_path(self.year_actual)
        self.forecasted_plan_path = get_forecasted_plan_path(self.year_actual)
        self.budget_path = get_budget_opex_path(self.year_actual)
        self.services_validated_paths = []

        self.view = None
        self.field_controller = None
        # -------------------------
        # 📊 Carga de gestores
        # -------------------------
        try:
            self.plan_actividades = PlanAnualActividades(self.data_loader, self.plan_path)
            self.opex_manager = OpexDataManager(self.data_loader, self.plan_path)
            self.opex_manager.load_opex_data()
            self.capacity_manager = OperativeCapacityManager(self.operative_capacity_file)
            self.cdf_df = self.data_loader.load_cdf_activities(self.data_loader, self.year_actual)
            print("✅ Gestores de datos de oficina cargados correctamente.")
        except (FileNotFoundError, pd.errors.EmptyDataError) as e:
            print(f"⚠️ ADVERTENCIA: No se encontraron o están vacíos los archivos de oficina: {e}.")
            print("    Las funcionalidades de oficina operarán con datos vacíos. Esto es normal para usuarios de campo.")
            QMessageBox.warning(None, "Archivos de Oficina no Encontrados",
                                f"No se pudieron cargar los archivos de configuración de oficina: {e}\n\n"
                                "Las funcionalidades de reportes de oficina estarán deshabilitadas o usarán datos vacíos.")
            self.plan_actividades = _EmptyPlanActividades()
            self.opex_manager = _EmptyOpexDataManager()
            self.capacity_manager = _EmptyCapacityManager()
            self.cdf_df = pd.DataFrame()

        # -------------------------
        # 📦 Configuración de reportes
        # -------------------------
        self.reports = [

            #{"class": RigReport, "title": "1.01 WI Rig", "type": "rig_schedule", "params": {"year": self.year_actual, "merged_opex_data": None, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": MISwacoReport, "title": "1.02 MI Swaco", "type": "mi_swaco", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": CompletionsReport, "title": "1.03 Completions", "type": "completions", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": BitsDrillingTRemedialReport, "title": "1.04 Bits, Drilling Tools & Remedial", "type": "bits_d_tools_remedial", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": SurfaceSystemsReport, "title": "1.05 Surface Systems", "type": "surface_systems", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": WirelineReport, "title": "1.06 Wireline Report", "type": "wireline", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": WellServicesReport, "title": "1.07 Well Services", "type": "well_services", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": TestingFluidAnalysisReport, "title": "1.08 Testing and Fluid Analysis", "type": "well_services", "params": {"year": self.year_actual, "operative_capacity": None,  "opex_manager": None, "plan_actividades": None}},
            #{"class": TubularsReport, "title": "1.09 Tubulars Report", "type": "tubulars", "params": {"year": self.year_actual, "operative_capacity": None, "plan_actividades": None, "opex_manager": None}},
            {"class": ServicesReport, "title": "1.10 Services", "type": "services", "params": {"year": self.year_actual, "operative_capacity": None, "plan_actividades": None, "opex_manager": None}},
            #{"class": EnvironmentReport, "title": "1.11 Environment", "type": "environment", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": ArtificialLiftReport, "title": "1.13 Artificial Lift", "type": "artificial_lift", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": IntegratedServicesReport, "title": "1.14 Integrates Services Management", "type": "well_services", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},
            #{"class": TanksAndTrunksReport, "title": "1.15 Tanks and Trunks", "type": "tanks_and_trunks", "params": {"year": self.year_actual, "operative_capacity": None, "opex_manager": None, "plan_actividades": None}},

        ]
        
    def set_view(self, view):
        """Asigna la vista principal al controlador e inicializa los sub-controladores."""
        self.view = view
        self.field_controller = FieldController(self.view)

    def generate_reports(self):
        activities_data = build_activities_dataframe(self.data_loader, self.plan_actividades, self.year_actual)

        for report_info in self.reports:
            report_class = report_info["class"]
            params = report_info["params"]

            # Agregar automáticamente opex_manager y plan_actividades si son requeridos
            if "opex_manager" in params:
                params["opex_manager"] = self.opex_manager
            if "plan_actividades" in params:
                params["plan_actividades"] = self.plan_actividades
            if "operative_capacity" in params:
                params["operative_capacity"] = self.capacity_manager.df

            report_instance = report_class(self.data_loader, **params)

            # Si el reporte es '1.10 Services', cargar duración del catálogo
            if report_info["title"] == "1.10 Services":
                try:
                    services_catalog_path = os.path.join(self.catalog_dir, "catalogo_solo_valores.xlsx")
                    df_services = pd.read_excel(services_catalog_path, sheet_name="Services")
                    inputs = self.get_services_inputs_from_table(df_services)
                    report_instance.set_manual_duration(inputs["duration"])
                    report_instance.set_manual_input_target_cost(inputs["target_cost"])
                    print(f"✅ Se aplicó duración manual desde catálogo: {inputs['duration']} días")
                except Exception as e:
                    print(f"⚠️ No se pudo aplicar duración desde catálogo para Services: {e}")
                
                if hasattr(self, "services_validated_paths"):
                    report_instance.set_validated_paths(self.services_validated_paths)

            forecast = report_instance.generate_forecast()
            budget = report_instance.generate_budget()
            deviations = report_instance.generate_deviations()
            graph = report_instance.generate_graph(forecast, budget, activities_data)            
            self.view.show_plot_view(
                graph,            deviations,
                title=report_info["title"],
                deviation_type=report_info.get("type", "default")
            )
    def open_office_activities_plan(self):
        """Abre la vista de planificación de actividades de oficina"""
        """"Abre la vista de planificación de actividades de oficina."""
        try:
            office_lines = [
                "1.06 Wireline Report",
                "1.08 Testing and Fluid Analysis",
                "1.15 Tanks and Trunks"
            ]
            save_path = get_planning_cost_path(self.year_actual)
            self.office_planning_view = OfficePlanningView(
                available_line_titles=office_lines,
                data_loader=self.data_loader,
                plan_actividades=self.plan_actividades,
                save_path=save_path,
                controller=self
            )
            self.office_planning_view.setWindowModality(Qt.ApplicationModal)
            self.office_planning_view.exec_()
        except FileNotFoundError as e:
            QMessageBox.warning(
                self.view,
                "File Not Found",
                "You must create the file first." + str(e)
            )

    def calcular_plan_costo_automatico_tanks_and_trunks(self, year, plan_actividades):
        """
        Calcula el costo automático para Tanks and Trunks y retorna el DataFrame con MONTH y PLANNED_COST.
        """
        report = TanksAndTrunksReport(
            data_loader=self.data_loader,
            year=year,
            operative_capacity=None,
            opex_manager=None,
            plan_actividades=plan_actividades
        )
        df_cost = report.generate_plan_cost_logic()
        return df_cost
    
    def calcular_plan_costo_automatico_testing(self, year, plan_actividades):
        report = TestingFluidAnalysisReport(
            data_loader=self.data_loader,
            year=year,
            operative_capacity=None,
            opex_manager=None,
            plan_actividades=plan_actividades
        )
        df_cost = report.generate_plan_cost_logic()
        return df_cost
    
    def calcular_plan_costo_automatico_wireline(self, year, plan_actividades):
        report = WirelineReport(
            data_loader=self.data_loader,
            year=year,
            operative_capacity=None,
            opex_manager=None,
            plan_actividades=plan_actividades
        )
        df_cost = report.generate_plan_cost_logic()
        return df_cost

    def regenerar_reporte_y_retorna_datos(self, title):
        """
        Regenera los datos de un solo reporte (general o línea de campo) y retorna:
        gráfico, desviaciones, comentarios (texto), tipo
        """

        self.plan_actividades = PlanAnualActividades(self.data_loader, self.plan_path)

        if self.field_controller and self.field_controller.is_field_report(title):
            return self.field_controller.regenerate_report_and_get_data(title)

        activities_data = build_activities_dataframe(self.data_loader, self.plan_actividades, self.year_actual)
        for report_info in self.reports:
            if report_info["title"] == title:
                report_class = report_info["class"]
                params = report_info["params"].copy()
                if "opex_manager" in params:
                    params["opex_manager"] = self.opex_manager
                if "plan_actividades" in params:
                    params["plan_actividades"] = self.plan_actividades
                if "operative_capacity" in params:
                    params["operative_capacity"] = self.capacity_manager.df
                instance = report_class(self.data_loader, **params)
                # Si es el reporte de 1.10 Services, aplicar duración desde el catálogo
                if title == "1.10 Services":
                    try:
                        services_catalog_path = os.path.join(self.catalog_dir, "catalogo_solo_valores.xlsx")
                        df_services = pd.read_excel(services_catalog_path, sheet_name="Services")
                        inputs = self.get_services_inputs_from_table(df_services)
                        duration = inputs["duration"] if inputs["duration"] > 0 else 1
                        instance.set_manual_duration(duration)
                        instance.set_manual_input_target_cost(inputs["target_cost"])
                        print(f"⏱ Duración manual aplicada a '1.10 Services': {duration} días")
                    except Exception as e:
                        print(f"⚠️ Error aplicando duración desde catálogo para '1.10 Services': {e}")
                    if hasattr(self, "services_validated_paths"):
                        instance.set_validated_paths(self.services_validated_paths)
                        
                forecast = instance.generate_forecast()
                budget = instance.generate_budget()
                deviations = instance.generate_deviations()
                graph = instance.generate_graph(forecast, budget, activities_data)
                # Comentarios (si está implementado)
                comentario = ""
                if hasattr(self, "get_comments_for_title"):
                    comentario = self.get_comments_for_title(title)
                return graph, deviations, comentario, report_info.get("type", "default")
        # Si no se encuentra el reporte
        return None, None, "", "default"

    def get_comments_for_title(self, title, line_type=None):
        if self.field_controller and self.field_controller.is_field_report(title):
            return self.field_controller.get_comments_for_title(title)
        else:
            mes_actual = datetime.now().strftime("%Y-%m")
            df = getattr(self, 'comments_df', None)
            if df is None or df.empty:
                return ""
            comentarios = df[(df["Report Title"] == title) & (df["Mes"] == mes_actual)]
            if comentarios.empty:
                return ""
            return "\n".join(f" {row['Comentario']}" for _, row in comentarios.iterrows())

    def save_new_comment(self, title, comentario):
        """Delega el guardado de comentarios al controlador correspondiente."""
        if self.field_controller and self.field_controller.is_field_report(title):
            self.field_controller.save_new_comment(title, comentario)
        else:
            now = datetime.now()
            mes_actual = now.strftime("%Y-%m")
            nuevo = pd.DataFrame([{
                "Report Title": title,
                "Mes": mes_actual,
                "Fecha": now.strftime("%Y-%m-%d %H:%M:%S"),
                "Usuario": getpass.getuser(),
                "Comentario": comentario
            }])
            
            df = self.comments_df
            if df is not None and not df.empty:
                df = df[~((df["Report Title"] == title) & (df["Mes"] == mes_actual))]
                self.comments_df = pd.concat([df, nuevo], ignore_index=True)
            else:
                self.comments_df = nuevo 
            save_comment(self.comments_df)
            print(f"💬 Comentario de oficina guardado para {title}")

    def open_well_selector_dialog(self):
        """
        Abre el selector de pozos para el reporte '1.10 Services'.

        Carga la lista de pozos disponibles y muestra una ventana
        para que el usuario seleccione manualmente los pozos que se usarán.
        La selección se guarda en un archivo Excel.
        """
        ruta_guardado = get_selected_services_wells_path()
        report = ServicesReport(
            self.data_loader,
            self.year_actual,
            self.capacity_manager.df,
            self.plan_actividades, 
            self.opex_manager
        )
        df_wells = report.load_available_wells()
        dialog = WellSelectorDialog(df_wells, ruta_guardado)
        dialog.exec_()

    def get_services_inputs_from_table(self, df_services):
        duration_row = df_services[(df_services["line"] == "__INPUT__") & (df_services["TIPO"] == "duration")]
        cost_row = df_services[(df_services["line"] == "__INPUT__") & (df_services["TIPO"] == "target_cost")]

        duration = float(duration_row["valor"].values[0]) if not duration_row.empty else 1
        target_cost = float(cost_row["valor"].values[0]) if not cost_row.empty else 0

        return {"duration": duration, "target_cost": target_cost}

    def initialize_operative_capacity(self):
        """
        Inicializa el gestor de capacidad operativa usando el archivo configurado.
        """
        self.capacity_manager = OperativeCapacityManager(self.operative_capacity_file)

    def open_services_forecast_path(self):
        """
        Abre la vista para gestionar la ruta de forecast de la línea 1.10 Services.
        Al cerrar la vista, recupera los registros validados y los pasa al modelo.
        """
        services_catalog_path = get_forecast_services_path_file()
        self.services_path_dialog = ServicesForecastPathView(services_catalog_path, controller=self)
        self.services_path_dialog.setWindowModality(Qt.NonModal)
        def on_paths_saved(registros_validados):
            self.services_validated_paths = registros_validados
            services_report = self.get_services_report_instance()
            services_report.set_validated_paths(registros_validados)
        self.services_path_dialog.paths_saved.connect(on_paths_saved)
        self.services_path_dialog.show()

    def open_avg_days_dialog(self):
        """
        Abre el diálogo para establecer el promedio de días.
        """
        dialog = AvgDaysDialog(controller=self)
        dialog.setWindowModality(Qt.ApplicationModal)
        dialog.exec_()

    def set_avg_days(self, value):
        """
        Recibe el valor desde AvgDaysDialog y lo envía al OperativeCapacityManager.
        """
        if hasattr(self, 'capacity_manager') and self.capacity_manager:
            self.capacity_manager.set_days_avg(value)
        else:
            print("No se pudo settear el promedio de días: capacity_manager no inicializado.")

    def get_all_values_to_show(self):
        try:
            # Crear una instancia del reporte de servicios
            services_report = self.get_services_report_instance()
            # Costo - Verificar los 3 tipos posibles

            target_cost = services_report.get_target_cost()
            calculated_cost = services_report.get_cost_per_day_automatically()
            total_cost = services_report.get_total_avg_cost()
            target_day = services_report.get_custom_duration()
            calculated_avg_days = services_report.get_avg_day_duration_automatically()

            return {
            "Costo Target": target_cost,
            "Costo Promedio": calculated_cost,
            "Costo Total": total_cost,
            "Días Target": target_day,
            "Días Promedio": calculated_avg_days
            }
            
        except Exception as e:
            print(f"No se pudo traer alguna variable calculada: {str(e)}")
            return None         

    def generate_forecast_by_path(self, id_costo, id_dia):
        """
        Genera el forecast usando el método generate_forecast_by_path del ServicesReport.
        
        :param id_costo: ID del método de cálculo de costo
        :param id_dia: ID del método de cálculo de días
        :return: DataFrame con el resultado del forecast
        """
        try:
            services_report = self.get_services_report_instance()
            result_df = services_report.generate_forecast_by_path(id_costo, id_dia)
            return result_df
            
        except Exception as e:
            print(f"Error en generate_forecast_by_path: {str(e)}")
            return None

    def get_services_report_instance(self):
        """
        Crea y retorna una instancia del reporte de servicios con los parámetros necesarios.
        """
        services_report = ServicesReport(
            self.data_loader,
            self.year_actual,
            self.capacity_manager.df,
            self.plan_actividades,
            self.opex_manager
        )
        services_catalog_path = os.path.join(self.catalog_dir, "catalogo_solo_valores.xlsx")
        df_services = pd.read_excel(services_catalog_path, sheet_name="Services")
        inputs = self.get_services_inputs_from_table(df_services)
        services_report.set_manual_duration(inputs["duration"])
        services_report.set_manual_input_target_cost(inputs["target_cost"])
        selected_path = get_selected_services_wells_path()
        services_report.load_selected_wells(selected_path)
        return services_report

    def open_services_summary(self):
        """
        Abre un resumen de costos y duración del reporte '1.10 Services' de forma no-modal.
        """
        try:
            
            services_catalog_path = os.path.join(self.catalog_dir, "catalogo_solo_valores.xlsx")
            df_services = pd.read_excel(services_catalog_path, sheet_name="Services")
            inputs = self.get_services_inputs_from_table(df_services)
            services_report = ServicesReport(
                self.data_loader,
                datetime.now().year,
                self.capacity_manager.df,
                self.plan_actividades,
                self.opex_manager
            )

            ruta_seleccionados = os.path.join(self.catalog_dir, "selected_services_wells.xlsx")
            services_report.load_selected_wells(ruta_seleccionados)
            df_budget = services_report.load_available_wells()
            df_durations = self.data_loader.calcular_duracion_promedio()
            costo_sel, costo_all, duracion_prom = services_report.get_costos_y_duracion(df_budget, df_durations)
            self.summary_dialog = ServicesResumenDialog(
                costo_target=inputs["target_cost"],
                costo_seleccion=costo_sel,
                costo_total=costo_all,
                duracion_target=inputs["duration"],
                duracion_promedio=duracion_prom
            )
            self.summary_dialog.show()
        except Exception as e:
            print(f"❌ Error al mostrar resumen de Services: {e}")
            # Opcional: Mostrar un QMessageBox con el error
            QMessageBox.critical(self.view, "Error", f"No se pudo generar el resumen: {e}")

    def open_tubulars_config(self):
        """
        Abre la ventana de configuración de tuberías (tubulars_config.xlsx),
        permitiendo múltiples filas por mes (Month, PipeDesc, Feet).
        """
        from views.tubulars_config import TubularsConfigDialog
        import pandas as pd
        from utils.file_manager import get_tubulars_config_path, get_catalog_path

        # 1️⃣ Ruta del config
        config_path = get_tubulars_config_path()

        # 2️⃣ Cargar df_config
        if os.path.exists(config_path):
            df_config = pd.read_excel(config_path)
        else:
            df_config = pd.DataFrame(columns=["Month", "PipeDesc", "Feet"])

        # 3️⃣ Cargar el catálogo unificado, filtrar line="Tubulars", cost_type="PER_FT"
        cat_path = get_catalog_path()
        df_catalog = self.data_loader.load_catalog_data(cat_path, sheet_name="Tubulars")

        df_tubing = df_catalog[
            (df_catalog["line"] == "Tubulars") & (df_catalog["cost_type"] == "PER_FT")
        ].copy()

        pipe_catalog = {
            str(row["description"]): float(row["cost_value"])
            for _, row in df_tubing.iterrows()
        }

        # 4️⃣ Crear el diálogo
        dlg = TubularsConfigDialog(df_config, pipe_catalog, parent=None)
        if dlg.exec_() == dlg.Accepted:
            df_updated = dlg.get_updated_df()
            df_updated.to_excel(config_path, index=False)
            print("✅ tubulars_config.xlsx actualizado.")

            
    def open_mi_swaco_config(self):
        """
        Abre la ventana de configuración del catálogo MI Swaco.
        (Versión inicial)
        """
        try:
            # Por ahora, solo abrimos el diálogo vacío
            dlg = MISwacoConfigDialog(parent=self.view)
            dlg.exec_()
        except Exception as e:
            QMessageBox.critical(
                self.view, 
                "Error", 
                f"No se pudo abrir la configuración de MI Swaco: {e}"
            )

    def open_completions_config(self):
        """
        Abre la ventana de configuración del catálogo Completions.
        (Versión inicial)
        """
        try:
            # Por ahora, solo abrimos el diálogo vacío
            dlg = CompletionsConfigDialog(parent=self.view)
            dlg.exec_()
        except Exception as e:
            QMessageBox.critical(
                self.view, 
                "Error", 
                f"No se pudo abrir la configuración de Completions: {e}"
            )

    def open_activity_plan_viewer(self):
        from views.plan_editor import PlanEditorWindow
        try:
            df_plan = self.plan_actividades.data_loader.load_plan_actividades_from_excel(
                self.plan_actividades.plan_path,
                self.plan_actividades.sheet_name
            )
            self.plan_editor = PlanEditorWindow(df_plan, self.plan_path)
            self.plan_editor.show()
        except Exception as e:
            print(f"❌ Error al abrir el archivo de plan: {e}")

    def open_forecasted_activity_plan_editor(self):
        from views.forecasted_plan_editor import ForecastedPlanEditorWindow
        sheet_name = f"ForecastedPlan{self.year_actual}"

        try:
            self.plan_actividades.plan_df = self.plan_actividades.data_loader.load_plan_actividades_from_excel(
                self.forecasted_plan_path,
                sheet_name
            )
            # 🟢 CAMBIO: Usar distribución híbrida que preserva valores guardados
            forecasted_plan_df = self.plan_actividades.calcular_distribucion_hibrida(
                year=self.year_actual, 
                saved_excel_path=self.forecasted_plan_path,
                saved_sheet_name=sheet_name
            )
            pozos_sugeridos = self.capacity_manager.get_total_tentative_opex_wells()
            
            self.plan_editor = ForecastedPlanEditorWindow(
                forecasted_plan_df, 
                self.forecasted_plan_path, 
                pozos_sugeridos,
                controller=self  
            )
            
            self.plan_editor.show()
        except Exception as e:
            print(f"❌ Error al abrir el editor de plan con CDF: {e}")

    def open_capex_config_editor(self):
        """
        Abre el editor de configuración de meses CAPEX. Carga los datos
        actuales, muestra el diálogo y guarda los cambios si el usuario acepta.
        """
        try:
            current_config = self.capex_config_service.load_config()
            dialog = CapexConfigDialog(current_config)
            
            if dialog.exec_() == QDialog.Accepted:
                updated_config = dialog.get_updated_config()
                self.capex_config_service.save_config(updated_config)
                QMessageBox.information(
                    None, 
                    "Configuración Guardada", 
                    "La configuración de meses CAPEX se ha guardado correctamente."
                )
        except Exception as e:
            QMessageBox.critical(
                None, 
                "Error", 
                f"No se pudo abrir o guardar la configuración CAPEX: {e}"
            )

    def cargar_plan_opex(self, file_path):
        self.opex_manager = OpexDataManager(file_path)
        return self.opex_manager.load_opex_plan(file_path)  # función que carga el plan
    
    def show_opex_editor(self):
        """
        Muestra la interfaz para editar manualmente el presupuesto OPEX.

        Carga los datos actuales de OPEX y abre la ventana de edición.
        """
        opex_df = self.opex_manager.load_opex_data()
        self.opex_editor = OpexEditorWindow(opex_df, self.save_edited_opex)
        self.opex_editor.show()


    def save_edited_opex(self, updated_df):
        """
        Guarda los cambios realizados en el presupuesto OPEX.

        Actualiza el gestor de datos y guarda los resultados en un archivo Excel.
        """
        self.opex_manager.set_opex_data(updated_df)
        self.opex_manager.save_opex_to_excel()
        print("✅ OPEX updated and saved to Excel.")

    def generate_all_slides(self, year_override=None, month_override=None):
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        activities_data = build_activities_dataframe(self.data_loader, self.plan_actividades, self.year_actual)

        for report_info in self.reports:
            report_class = report_info["class"]
            params = report_info["params"].copy()

            # Actualizar dependencias
            if "opex_manager" in params:
                params["opex_manager"] = self.opex_manager
            if "plan_actividades" in params:
                params["plan_actividades"] = self.plan_actividades
            if "operative_capacity" in params:
                params["operative_capacity"] = self.capacity_manager.df

            instance = report_class(self.data_loader, **params)

            forecast = instance.generate_forecast()
            budget = instance.generate_budget()
            deviations_df = instance.generate_deviations()
            graph = instance.generate_graph(forecast, budget, activities_data)

            comments = self.view.get_comments_for_title(report_info["title"])
            if report_info["title"] == "1.13 Artificial Lift":
                # Obtener el primer día del mes actual y restar un día => último día del mes anterior
                last_month_date = datetime.now().replace(day=1) - timedelta(days=1)
                closing_month = last_month_date.strftime("%B")  # e.g., "March"# e.g., "April"
                deviations_str = "\n".join(self.format_closing_month_artificial_lift_deviations(deviations_df, closing_month))
            else:
                deviations_str = "\n".join([str(row.to_dict()) for _, row in deviations_df.iterrows()]) if not deviations_df.empty else "No deviations found."

            add_slide_to_presentation(prs, graph, deviations_str, comments, title=report_info["title"])

        output_path = get_output_path_for_pptx(year=year_override, month=month_override)
        prs.save(output_path)
        print(f"✅ Presentación exportada en: {output_path}")

    def format_closing_month_artificial_lift_deviations(self, deviations, closing_month):
        from utils.dates import normalize_month_names

        deviation_groups = [
            'Servicio_Deviation', 'Equipo_Deviation', 'Protectores de Cable_Deviation',
            'Capilar_Deviation', 'Equipo Superficie_Deviation', 'Desarenador_Deviation',
            'Cable Nuevo_Deviation', 'B&H_Deviation'
        ]
        threshold = 20000

        deviations = deviations.copy()
        deviations["MONTH"] = normalize_month_names(deviations["MONTH"])
        month_order = [
            "January", "February", "March", "April", "May", "June",
            "July", "August", "September", "October", "November", "December"
        ]
        deviations["MONTH"] = pd.Categorical(deviations["MONTH"], categories=month_order, ordered=True)
        deviations = deviations[deviations["MONTH"] == closing_month].sort_values(by=["MONTH", "WELL"])

        formatted = []
        for _, row in deviations.iterrows():
            line = f"{row['WELL']} ({row['MONTH']} {int(row['YEAR'])}): "
            details = [
                f"{group.replace('_Deviation', '')}: {row[group]:+,.2f} $" 
                for group in deviation_groups 
                if abs(row.get(group, 0)) > threshold
            ]
            if details:
                line += " | ".join(details)
                formatted.append(line)
        return formatted
    
    def update_capacity_value(self, row_index, column, value):
        """
        Actualiza un valor en la tabla de capacidad operativa.
        """
        try:
            self.capacity_manager.update_value(row_index, column, value)
            self.dataUpdated.emit()
        except (IndexError, ValueError) as e:
            print(f"Error actualizando capacidad operativa: {e}")

    def save_operative_capacity(self, file_path=None):
        """
        Guarda los datos de capacidad operativa a Excel.
        """
        self.capacity_manager.save(file_path)

    def open_table_popup(self):
        """
        Muestra la tabla de capacidad operativa en la vista.
        """
        self.view.show_table(self.capacity_manager.df)

    def save_table_data_to_excel(self, file_path):
        """
        Guarda una copia del archivo de capacidad operativa.
        """
        if self.capacity_manager.df.empty:
            print("No hay datos para guardar.")
            return

        self.capacity_manager.export_to(file_path)

    def update_capex_from_cdf(self):
        """
        Llama al DataLoader para obtener y distribuir días de CAPEX por mes,
        y actualiza self.data.
        """
        capex_monthly = self.data_loader.fetch_and_distribute_capex(year=2025)  # o None para el actual
        if not capex_monthly:
            print("No se pudo actualizar CAPEX desde CDF (diccionario vacío).")
            return
        # export capex monthly to exel
        capex_df = pd.DataFrame(capex_monthly.items(), columns=["Mes", "Días CAPEX"])

        # Diccionario de mapeo de número de mes a nombre en inglés
        month_map = {
            1: "January",
            2: "February",
            3: "March",
            4: "April",
            5: "May",
            6: "June",
            7: "July",
            8: "August",
            9: "September",
            10: "October",
            11: "November",
            12: "December"
        }

        for row in self.data:
            mes_num = row["Mes"]  # aquí viene el mes como número (p.ej. 1, 2, 3...)
            mes_en = month_map.get(mes_num, None)  # Convertimos a "January", "February", etc.
            
            # Si existe el mes en el mapeo, entonces asignamos el valor de capex_monthly
            if mes_en:
                row["Días CAPEX"] = capex_monthly.get(mes_en, 0)
            else:
                row["Días CAPEX"] = 0
        self.recalculate_table_data()
        self.dataUpdated.emit()
        print("Días CAPEX actualizados correctamente.")

    # Refactor para TODO LO DE CAMPO (Main Controller se queda como fachada)
    def open_initial_activities_plan(self):
        """
        Este método recibe la llamada desde la vista principal y la delega al sub-controlador de campo correspondiente.
        """
        if self.field_controller:
            self.field_controller.open_adaptive_planning()
        else:
            print("Error: FieldController no ha sido inicializado.")
    
    def generate_field_reports(self):
        """Delega la generación de reportes de campo al FieldController."""
        if self.field_controller:
            self.field_controller.generate_field_reports()
    
    def open_approved_budget_and_activities(self):
        """Delega la apertura del presupuesto aprobado al FieldController."""
        if self.field_controller:
            self.field_controller.open_approved_budget_and_activities()

    def open_leader_summary_report(self):
        if self.field_controller:
            self.field_controller.open_leader_summary_report()

    def open_south_zone_quote_extractor_view(self):
         if self.field_controller:
            self.field_controller.open_south_zone_quote_extractor_view()

    def open_quote_extractor_view(self):
        if self.field_controller:
            self.field_controller.open_quote_extractor_view()

    def open_historical_initial_cost_view(self):
        if self.field_controller:
            self.field_controller.open_historical_initial_cost_view()

    def open_total_executed_activities_view(self, selected_line):
        if self.field_controller:
            self.field_controller.open_total_executed_activities_view(selected_line)

    def open_categorizer_executed_activities_catalog(self, selected_line):
        if self.field_controller:
            self.field_controller.open_categorizer_executed_activities_catalog(selected_line)
    
    def open_cpi_spi_by_line(self):
        if self.field_controller:
            self.field_controller.open_cpi_spi_by_line()

    def regenerar_reporte_linea_campo(self, line_title):
        if self.field_controller:
            self.field_controller.regenerar_reporte_linea_campo(line_title)

    def get_monthly_summary_dataframe(self):
        if self.field_controller:
            return self.field_controller.generate_leader_line_report()
        return pd.DataFrame()
    
    def generate_leader_line_report(self):
        """Delega la generación del reporte líder visual al FieldController."""
        if self.field_controller:
            self.field_controller.generate_leader_line_report()
