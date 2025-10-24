from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def add_slide_to_presentation(prs, graph, deviations_text, comments_text, title='1.01 Rig'):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Título
    title_box = slide.shapes.title
    title_box.text = title
    title_frame = title_box.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.name = "Arial Narrow"
    title_paragraph.font.color.rgb = RGBColor(0x00, 0x70, 0x64)

    # Gráfico
    img_stream = io.BytesIO()
    graph.savefig(img_stream, format='png')
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(10.5))

    # Deviations
    deviations_box = slide.shapes.add_textbox(Inches(11.2), Inches(1.5), Inches(4), Inches(3))
    frame = deviations_box.text_frame
    p = frame.add_paragraph()
    p.text = "Deviations:\n" + deviations_text
    p.font.size = Pt(15)
    p.font.name = "Arial Narrow"
    frame.word_wrap = True

    # Comments
    comments_box = slide.shapes.add_textbox(Inches(11.2), Inches(4.5), Inches(4), Inches(2))
    frame = comments_box.text_frame
    p = frame.add_paragraph()
    p.text = "Comments:\n" + comments_text
    p.font.size = Pt(15)
    p.font.name = "Arial Narrow"
    frame.word_wrap = True
